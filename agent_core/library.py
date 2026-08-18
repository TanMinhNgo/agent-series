"""Personal local file library used by uploads and agent-created artifacts."""

from __future__ import annotations

import csv
import json
from io import BytesIO, StringIO
from pathlib import Path
from uuid import uuid4

from sqlalchemy import func, select

from .storage import Database, LibraryAsset

ALLOWED_SUFFIXES = {".pdf", ".docx", ".xlsx", ".pptx", ".md", ".csv", ".json", ".txt", ".png", ".jpg", ".jpeg", ".webp", ".gif"}
MAX_FILE_BYTES = 25 * 1024 * 1024


class LibraryService:
    def __init__(self, database: Database, directory: Path):
        self.database, self.directory = database, directory
        self.directory.mkdir(parents=True, exist_ok=True)

    def list(self, query: str = "", project_id: str | None = None, scope: str = "all") -> list[LibraryAsset]:
        with self.database.session() as session:
            latest = select(LibraryAsset.artifact_id, func.max(LibraryAsset.version).label("version")).group_by(LibraryAsset.artifact_id).subquery()
            statement = select(LibraryAsset).join(latest, (LibraryAsset.artifact_id == latest.c.artifact_id) & (LibraryAsset.version == latest.c.version)).order_by(LibraryAsset.created_at.desc())
            if query.strip():
                statement = statement.where(LibraryAsset.name.ilike(f"%{query.strip()}%"))
            if scope == "global":
                statement = statement.where(LibraryAsset.project_id.is_(None))
            elif scope == "project" and project_id:
                statement = statement.where(LibraryAsset.project_id == project_id)
            return list(session.scalars(statement))

    def upload(
        self,
        name: str,
        content_type: str,
        data: bytes,
        source: str = "upload",
        project_id: str | None = None,
    ) -> LibraryAsset:
        suffix = Path(name).suffix.lower()
        if suffix not in ALLOWED_SUFFIXES: raise ValueError("Định dạng file chưa được hỗ trợ trong Thư viện.")
        if not data or len(data) > MAX_FILE_BYTES: raise ValueError("File phải có dung lượng từ 1 byte đến 25 MB.")
        stored_name = f"library_{uuid4().hex}{suffix}"
        (self.directory / stored_name).write_bytes(data)
        with self.database.session() as session:
            asset = LibraryAsset(
                name=Path(name).name[:255],
                stored_name=stored_name,
                mime_type=content_type or "application/octet-stream",
                size_bytes=len(data),
                source=source,
                project_id=project_id,
                is_project_source=bool(project_id),
                index_status="queued" if project_id else "pending",
            )
            session.add(asset); session.commit(); return asset

    def versions(self, asset_id: str) -> list[LibraryAsset]:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                return []
            return list(session.scalars(select(LibraryAsset).where(LibraryAsset.artifact_id == asset.artifact_id).order_by(LibraryAsset.version.desc())))

    def update(self, asset_id: str, *, name: str | None = None, project_id: str | None = None, is_project_source: bool | None = None) -> LibraryAsset | None:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                return None
            versions = session.scalars(select(LibraryAsset).where(LibraryAsset.artifact_id == asset.artifact_id)).all()
            if name is not None:
                clean_name = Path(name).name.strip()[:255]
                if not clean_name:
                    raise ValueError("Tên artifact không được để trống.")
                for item in versions:
                    item.name = clean_name
            if project_id is not None or is_project_source is False:
                for item in versions:
                    item.project_id = project_id
                    item.is_project_source = bool(project_id) if is_project_source is None else bool(is_project_source and project_id)
                    if item.is_project_source and item.index_status == "pending":
                        item.index_status = "queued"
            elif is_project_source is True:
                if not asset.project_id:
                    raise ValueError("Cần chọn Project trước khi ghim artifact.")
                for item in versions:
                    item.is_project_source = True
                    if item.index_status == "pending":
                        item.index_status = "queued"
            session.commit()
            return session.get(LibraryAsset, asset_id)

    def create_version(self, asset_id: str, name: str, content_type: str, data: bytes) -> LibraryAsset:
        suffix = Path(name).suffix.lower()
        if suffix not in ALLOWED_SUFFIXES:
            raise ValueError("Định dạng file chưa được hỗ trợ trong Thư viện.")
        if not data or len(data) > MAX_FILE_BYTES:
            raise ValueError("File phải có dung lượng từ 1 byte đến 25 MB.")
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                raise ValueError("Không tìm thấy artifact.")
            version = (session.scalar(select(func.max(LibraryAsset.version)).where(LibraryAsset.artifact_id == asset.artifact_id)) or 0) + 1
            stored_name = f"library_{uuid4().hex}{suffix}"
            (self.directory / stored_name).write_bytes(data)
            item = LibraryAsset(
                name=asset.name,
                stored_name=stored_name,
                mime_type=content_type or "application/octet-stream",
                size_bytes=len(data),
                source="version",
                project_id=asset.project_id,
                artifact_id=asset.artifact_id,
                version=version,
                is_project_source=asset.is_project_source,
                index_status="queued" if asset.is_project_source else "pending",
            )
            session.add(item)
            session.commit()
            return item

    def delete(self, asset_id: str) -> bool:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None: return False
            path = self.directory / asset.stored_name
            if path.exists(): path.unlink()
            session.delete(asset); session.commit(); return True

    def create_export(self, name: str, format: str, content: str, project_id: str | None = None) -> LibraryAsset:
        format = format.lower().lstrip(".")
        if format not in {"docx", "xlsx", "pptx", "md", "csv", "pdf", "json"}: raise ValueError("Định dạng export chưa hỗ trợ.")
        filename = f"{Path(name).stem or 'tai-lieu'}.{format}"
        if format in {"md", "csv", "json"}:
            payload = content.encode("utf-8")
            mime = {"md": "text/markdown", "csv": "text/csv", "json": "application/json"}[format]
        elif format == "docx":
            from docx import Document
            doc = Document(); doc.add_paragraph(content); stream = BytesIO(); doc.save(stream); payload, mime = stream.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif format == "xlsx":
            from openpyxl import Workbook
            book = Workbook(); sheet = book.active
            for row in csv.reader(StringIO(content)): sheet.append(row)
            stream = BytesIO(); book.save(stream); payload, mime = stream.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif format == "pptx":
            from pptx import Presentation
            deck = Presentation()
            for slide_text in content.split("\n\n"):
                slide = deck.slides.add_slide(deck.slide_layouts[1]); lines = slide_text.splitlines(); slide.shapes.title.text = lines[0] if lines else "Nội dung"; slide.placeholders[1].text = "\n".join(lines[1:])
            stream = BytesIO(); deck.save(stream); payload, mime = stream.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            from reportlab.pdfgen.canvas import Canvas
            stream = BytesIO(); canvas = Canvas(stream); text = canvas.beginText(48, 800)
            for line in content.splitlines(): text.textLine(line[:120])
            canvas.drawText(text); canvas.save(); payload, mime = stream.getvalue(), "application/pdf"
        return self.upload(filename, mime, payload, source="generated", project_id=project_id)
