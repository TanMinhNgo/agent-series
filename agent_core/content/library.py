"""Personal local file library used by uploads and agent-created artifacts."""

from __future__ import annotations

import csv
import json
from io import BytesIO, StringIO
from pathlib import Path
from uuid import uuid4

from sqlalchemy import func, select

from ..persistence.store import Database, LibraryAsset, current_user_id
from .file_storage import FileStorageService

ALLOWED_SUFFIXES = {
    ".pdf", ".docx", ".xlsx", ".pptx", ".md", ".csv", ".json", ".txt", ".py", ".ts", ".tsx",
    ".png", ".jpg", ".jpeg", ".webp", ".gif",
}
MAX_FILE_BYTES = 25 * 1024 * 1024
EXPORT_FORMATS = {"docx", "xlsx", "pptx", "md", "csv", "pdf", "json", "txt", "py", "ts", "tsx"}
TEXT_EXPORT_MIME_TYPES = {
    "md": "text/markdown", "csv": "text/csv", "json": "application/json", "txt": "text/plain",
    "py": "text/x-python", "ts": "text/typescript", "tsx": "text/tsx",
}


class LibraryService:
    def __init__(self, database: Database, directory: Path, storage: FileStorageService | None = None):
        self.database, self.directory = database, directory
        self.directory.mkdir(parents=True, exist_ok=True)
        self.storage = storage or FileStorageService(directory)

    def list(self, query: str = "", project_id: str | None = None, scope: str = "all") -> list[LibraryAsset]:
        with self.database.session() as session:
            latest = select(LibraryAsset.artifact_id, func.max(LibraryAsset.version).label("version")).group_by(LibraryAsset.artifact_id).subquery()
            statement = select(LibraryAsset).join(latest, (LibraryAsset.artifact_id == latest.c.artifact_id) & (LibraryAsset.version == latest.c.version)).order_by(LibraryAsset.created_at.desc())
            if query.strip():
                statement = statement.where(LibraryAsset.name.ilike(f"%{query.strip()}%"))
            if scope == "global":
                statement = statement.where(LibraryAsset.project_id.is_(None))
            elif scope == "project" and project_id:
                statement = statement.where(LibraryAsset.project_id == project_id)
            return list(session.scalars(statement))

    def upload(
        self,
        name: str,
        content_type: str,
        data: bytes,
        source: str = "upload",
        project_id: str | None = None,
    ) -> LibraryAsset:
        suffix = Path(name).suffix.lower()
        if suffix not in ALLOWED_SUFFIXES: raise ValueError("Định dạng file chưa được hỗ trợ trong Thư viện.")
        if not data or len(data) > MAX_FILE_BYTES: raise ValueError("File phải có dung lượng từ 1 byte đến 25 MB.")
        stored = self.storage.upload(data, name, "library")
        with self.database.session() as session:
            asset = LibraryAsset(
                name=Path(name).name[:255],
                stored_name=stored.stored_name,
                storage_provider=stored.provider,
                storage_file_id=stored.file_id,
                mime_type=content_type or "application/octet-stream",
                size_bytes=len(data),
                source=source,
                project_id=project_id,
                is_project_source=bool(project_id),
                index_status="queued" if project_id else "pending",
            )
            session.add(asset); session.commit(); return asset

    def versions(self, asset_id: str) -> list[LibraryAsset]:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                return []
            return list(session.scalars(select(LibraryAsset).where(LibraryAsset.artifact_id == asset.artifact_id).order_by(LibraryAsset.version.desc())))

    def ensure_remote(self, asset_id: str) -> LibraryAsset | None:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None or asset.storage_provider != "local":
                return asset
            stored = self.storage.migrate_local(asset.stored_name, asset.name, "library")
            if stored is None:
                return asset
            asset.storage_provider, asset.stored_name, asset.storage_file_id = stored.provider, stored.stored_name, stored.file_id
            session.commit()
            return asset

    def update(self, asset_id: str, *, name: str | None = None, project_id: str | None = None, is_project_source: bool | None = None) -> LibraryAsset | None:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                return None
            versions = session.scalars(select(LibraryAsset).where(LibraryAsset.artifact_id == asset.artifact_id)).all()
            self._rename_versions(versions, name)
            self._update_project_source(asset, versions, project_id, is_project_source)
            session.commit()
            return session.get(LibraryAsset, asset_id)

    @staticmethod
    def _rename_versions(versions: list[LibraryAsset], name: str | None) -> None:
        if name is None:
            return
        clean_name = Path(name).name.strip()[:255]
        if not clean_name:
            raise ValueError("Tên artifact không được để trống.")
        for item in versions:
            item.name = clean_name

    @staticmethod
    def _queue_index_if_project_source(asset: LibraryAsset) -> None:
        if asset.is_project_source and asset.index_status == "pending":
            asset.index_status = "queued"

    def _update_project_source(self, asset: LibraryAsset, versions: list[LibraryAsset], project_id: str | None, is_project_source: bool | None) -> None:
        if project_id is not None or is_project_source is False:
            self._assign_project_source(versions, project_id, is_project_source)
        elif is_project_source is True:
            self._pin_existing_project_source(asset, versions)

    def _assign_project_source(self, versions: list[LibraryAsset], project_id: str | None, is_project_source: bool | None) -> None:
        for item in versions:
            item.project_id = project_id
            item.is_project_source = bool(project_id) if is_project_source is None else bool(is_project_source and project_id)
            self._queue_index_if_project_source(item)

    def _pin_existing_project_source(self, asset: LibraryAsset, versions: list[LibraryAsset]) -> None:
        if not asset.project_id:
            raise ValueError("Cần chọn Project trước khi ghim artifact.")
        for item in versions:
            item.is_project_source = True
            self._queue_index_if_project_source(item)

    def create_version(self, asset_id: str, name: str, content_type: str, data: bytes) -> LibraryAsset:
        suffix = Path(name).suffix.lower()
        if suffix not in ALLOWED_SUFFIXES:
            raise ValueError("Định dạng file chưa được hỗ trợ trong Thư viện.")
        if not data or len(data) > MAX_FILE_BYTES:
            raise ValueError("File phải có dung lượng từ 1 byte đến 25 MB.")
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                raise ValueError("Không tìm thấy artifact.")
            version = (session.scalar(select(func.max(LibraryAsset.version)).where(LibraryAsset.artifact_id == asset.artifact_id)) or 0) + 1
            stored = self.storage.upload(data, name, "library")
            item = LibraryAsset(
                name=asset.name,
                stored_name=stored.stored_name,
                storage_provider=stored.provider,
                storage_file_id=stored.file_id,
                mime_type=content_type or "application/octet-stream",
                size_bytes=len(data),
                source="version",
                project_id=asset.project_id,
                artifact_id=asset.artifact_id,
                version=version,
                is_project_source=asset.is_project_source,
                index_status="queued" if asset.is_project_source else "pending",
            )
            session.add(item)
            session.commit()
            return item

    def restore_version(self, asset_id: str) -> LibraryAsset:
        """Copy a selected historical version into a new, latest version."""
        asset = self.ensure_remote(asset_id)
        if asset is None:
            raise ValueError("Không tìm thấy artifact.")
        data = self.storage.read(asset.storage_provider, asset.stored_name, asset.storage_file_id)
        return self.create_version(asset.id, asset.name, asset.mime_type, data)

    def delete(self, asset_id: str) -> bool:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None: return False
            self.storage.delete(asset.storage_provider, asset.stored_name, asset.storage_file_id)
            session.delete(asset); session.commit(); return True

    def create_export(self, name: str, format: str, content: str, project_id: str | None = None) -> LibraryAsset:
        format = format.lower().lstrip(".")
        if format not in EXPORT_FORMATS: raise ValueError("Định dạng export chưa hỗ trợ.")
        filename = f"{Path(name).stem or 'tai-lieu'}.{format}"
        payload, mime = self._export_payload(format, content)
        return self.upload(filename, mime, payload, source="generated", project_id=project_id)

    @staticmethod
    def _export_payload(format: str, content: str) -> tuple[bytes, str]:
        if format in TEXT_EXPORT_MIME_TYPES:
            return content.encode("utf-8"), TEXT_EXPORT_MIME_TYPES[format]
        if format == "docx":
            from docx import Document
            doc = Document(); doc.add_paragraph(content); stream = BytesIO(); doc.save(stream)
            return stream.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if format == "xlsx":
            from openpyxl import Workbook
            book = Workbook(); sheet = book.active
            for row in csv.reader(StringIO(content)): sheet.append(row)
            stream = BytesIO(); book.save(stream)
            return stream.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        if format == "pptx":
            from pptx import Presentation
            deck = Presentation()
            for slide_text in content.split("\n\n"):
                slide = deck.slides.add_slide(deck.slide_layouts[1]); lines = slide_text.splitlines(); slide.shapes.title.text = lines[0] if lines else "Nội dung"; slide.placeholders[1].text = "\n".join(lines[1:])
            stream = BytesIO(); deck.save(stream)
            return stream.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        return LibraryService._pdf_export(content)

    @staticmethod
    def _pdf_export(content: str) -> tuple[bytes, str]:
        from reportlab.pdfgen.canvas import Canvas
        stream = BytesIO(); canvas = Canvas(stream); text = canvas.beginText(48, 800)
        for line in content.splitlines(): text.textLine(line[:120])
        canvas.drawText(text); canvas.save()
        return stream.getvalue(), "application/pdf"
