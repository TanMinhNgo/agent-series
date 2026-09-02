"""Versioned Project artifacts with preview extraction and pgvector retrieval."""

from __future__ import annotations

import csv
from dataclasses import dataclass
from difflib import unified_diff
import json
from io import BytesIO, StringIO
from pathlib import Path

from pypdf import PdfReader
from sqlalchemy import select

from ..knowledge.rag import _chunks
from ..persistence.store import ArtifactChunk, Database, LibraryAsset, current_user_id
from .file_storage import FileStorageService
from ..tools.base import ToolSpec

PREVIEW_LIMIT = 30_000
ARTIFACT_NOT_FOUND = "Không tìm thấy artifact."
EDITABLE_ARTIFACT_SUFFIXES = {".md", ".txt", ".json", ".py", ".ts", ".tsx"}


@dataclass(frozen=True)
class ArtifactEditContext:
    asset_id: str
    artifact_id: str
    name: str
    mime_type: str
    project_id: str | None
    version: int
    content: str


def extract_artifact_text(data: bytes | Path, suffix: str) -> str:
    """Extract safe plain text from the document types the Library can create."""
    if isinstance(data, Path):
        data = data.read_bytes()
    suffix = suffix.lower()
    if suffix in {".md", ".txt", ".csv", ".json", ".py", ".ts", ".tsx"}:
        return data.decode("utf-8", errors="replace")
    if suffix == ".pdf":
        return "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(data)).pages)
    if suffix == ".docx":
        from docx import Document
        return "\n".join(item.text for item in Document(BytesIO(data)).paragraphs)
    if suffix == ".xlsx":
        from openpyxl import load_workbook
        book = load_workbook(BytesIO(data), read_only=True, data_only=True)
        rows = []
        for sheet in book.worksheets:
            rows.append(f"# {sheet.title}")
            rows.extend("\t".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
        return "\n".join(rows)
    if suffix == ".pptx":
        from pptx import Presentation
        presentation = Presentation(BytesIO(data))
        return "\n".join(
            text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text") for text in [shape.text]
        )
    raise ValueError("File này chưa hỗ trợ trích xuất nội dung.")


class ArtifactService:
    def __init__(self, database: Database, directory: Path, embedding_model: str, storage: FileStorageService | None = None):
        self.database, self.directory, self.embedding_model_name = database, directory, embedding_model
        self.storage = storage or FileStorageService(directory)
        self._embedder = None

    def _embed(self, values: list[str], prefix: str) -> list[list[float]]:
        if self._embedder is None:
            from sentence_transformers import SentenceTransformer
            self._embedder = SentenceTransformer(self.embedding_model_name)
        return self._embedder.encode([f"{prefix}: {value}" for value in values], normalize_embeddings=True).tolist()

    def preview(self, asset_id: str) -> dict:
        asset = self._ensure_remote(asset_id)
        if asset is None:
            raise ValueError(ARTIFACT_NOT_FOUND)
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                raise ValueError(ARTIFACT_NOT_FOUND)
            suffix = Path(asset.name).suffix.lower()
            if suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
                return {"kind": "image"}
            if suffix == ".pdf":
                return {"kind": "pdf"}
            try:
                content = extract_artifact_text(self.storage.read(asset.storage_provider, asset.stored_name, asset.storage_file_id), suffix)
            except ValueError:
                return {"kind": "download"}
            return {"kind": "text", "content": content[:PREVIEW_LIMIT], "truncated": len(content) > PREVIEW_LIMIT}

    def edit_context(self, asset_id: str, project_id: str | None) -> ArtifactEditContext:
        """Read one exact, user-scoped artifact version for an AI edit request."""
        self._ensure_remote(asset_id)
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                raise ValueError(ARTIFACT_NOT_FOUND)
            if asset.project_id != project_id:
                raise ValueError("File này không thuộc Project của chat hiện tại.")
            suffix = Path(asset.name).suffix.lower()
            if suffix not in EDITABLE_ARTIFACT_SUFFIXES:
                raise ValueError("Chỉ hỗ trợ sửa bằng AI cho Markdown, text, JSON, Python và TypeScript.")
            content = extract_artifact_text(
                self.storage.read(asset.storage_provider, asset.stored_name, asset.storage_file_id), suffix,
            )
            if len(content) > PREVIEW_LIMIT:
                raise ValueError("File vượt quá 30.000 ký tự nên chưa thể sửa an toàn bằng AI.")
            return ArtifactEditContext(
                asset_id=asset.id,
                artifact_id=asset.artifact_id,
                name=asset.name,
                mime_type=asset.mime_type,
                project_id=asset.project_id,
                version=asset.version,
                content=content,
            )

    def diff(self, asset_id: str) -> dict:
        """Return a unified diff against the immediately preceding version."""
        self._ensure_remote(asset_id)
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                raise ValueError(ARTIFACT_NOT_FOUND)
            previous = session.scalar(
                select(LibraryAsset)
                .where(LibraryAsset.artifact_id == asset.artifact_id, LibraryAsset.version < asset.version)
                .order_by(LibraryAsset.version.desc())
                .limit(1)
            )
            if previous is None:
                raise ValueError("File này chưa có phiên bản trước để so sánh.")
            suffix = Path(asset.name).suffix.lower()
            if suffix not in EDITABLE_ARTIFACT_SUFFIXES:
                raise ValueError("Chỉ hỗ trợ xem thay đổi cho file text hoặc code.")
            current_content = extract_artifact_text(
                self.storage.read(asset.storage_provider, asset.stored_name, asset.storage_file_id), suffix,
            )
            previous_content = extract_artifact_text(
                self.storage.read(previous.storage_provider, previous.stored_name, previous.storage_file_id), suffix,
            )
            if max(len(current_content), len(previous_content)) > PREVIEW_LIMIT:
                raise ValueError("File vượt quá 30.000 ký tự nên chưa thể so sánh an toàn.")
            diff = "\n".join(
                unified_diff(
                    previous_content.splitlines(),
                    current_content.splitlines(),
                    fromfile=f"{previous.name} (v{previous.version})",
                    tofile=f"{asset.name} (v{asset.version})",
                    lineterm="",
                )
            )
            return {
                "baseAssetId": previous.id,
                "baseVersion": previous.version,
                "assetId": asset.id,
                "version": asset.version,
                "diff": diff,
            }

    def index(self, asset_id: str) -> LibraryAsset:
        self._ensure_remote(asset_id)
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None:
                raise ValueError(ARTIFACT_NOT_FOUND)
            asset.index_status, asset.index_error = "indexing", None
            session.query(ArtifactChunk).filter(ArtifactChunk.asset_id == asset_id).delete()
            session.commit()
            try:
                content = extract_artifact_text(self.storage.read(asset.storage_provider, asset.stored_name, asset.storage_file_id), Path(asset.name).suffix)
                parts = _chunks(content)
                if not parts:
                    asset.index_status, asset.index_error = "failed", "Artifact không có nội dung văn bản để truy hồi."
                    session.commit()
                    raise RuntimeError(asset.index_error)
                vectors = self._embed(parts, "passage")
                session.add_all(ArtifactChunk(asset_id=asset.id, chunk_index=index, content=part, embedding=vector) for index, (part, vector) in enumerate(zip(parts, vectors)))
                asset.index_status = "ready"
            except Exception as exc:  # noqa: BLE001
                asset.index_status, asset.index_error = "failed", str(exc)
            session.commit()
            if asset.index_status == "failed":
                raise RuntimeError(asset.index_error or "Không thể index artifact.")
            return asset

    def _ensure_remote(self, asset_id: str) -> LibraryAsset | None:
        with self.database.session() as session:
            asset = session.get(LibraryAsset, asset_id)
            if asset is None or asset.storage_provider != "local":
                return asset
            stored = self.storage.migrate_local(asset.stored_name, asset.name, "library")
            if stored is None:
                return asset
            asset.storage_provider, asset.stored_name, asset.storage_file_id = stored.provider, stored.stored_name, stored.file_id
            session.commit()
            return asset

    def search(self, query: str, project_id: str, top_k: int = 4) -> str:
        vector = self._embed([query], "query")[0]
        with self.database.session() as session:
            distance = ArtifactChunk.embedding.cosine_distance(vector)
            rows = session.execute(
                select(ArtifactChunk, LibraryAsset.name, LibraryAsset.version, distance.label("distance"))
                .join(LibraryAsset)
                .where(LibraryAsset.project_id == project_id, LibraryAsset.is_project_source.is_(True), LibraryAsset.index_status == "ready")
                .order_by(distance)
                .limit(max(1, min(top_k, 8)))
            ).all()
        if not rows:
            return "Không có Project Source phù hợp đã index."
        return "\n\n".join(
            f"[Artifact {number}: {name}, v{version}]\n{chunk.content}"
            for number, (chunk, name, version, _) in enumerate(rows, start=1)
        )


def build_artifact_tool(service: ArtifactService, project_id: str) -> ToolSpec:
    return ToolSpec(
        name="search_project_sources",
        description="Tìm trong artifact đã ghim vào Project. Khi dùng, nêu tên artifact và version làm nguồn.",
        parameters={"type": "object", "properties": {"query": {"type": "string"}, "top_k": {"type": "integer"}}, "required": ["query"]},
        func=lambda query, top_k=4: service.search(query, project_id, top_k),
    )
